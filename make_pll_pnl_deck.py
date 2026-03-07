#!/usr/bin/env python3
"""Generate an executive PPTX deck from pll_pnl_2025_vs_2026_report.md.

Output:
  /data/workspace/PLL_PnL_2025A_vs_2026B_Exec_Review.pptx

Requires:
  python-pptx, pandas (optional), numpy (optional)
"""

from __future__ import annotations

from pathlib import Path
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPORT_PATH = Path("/data/workspace/pll_pnl_2025_vs_2026_report.md")
OUT_PATH = Path("/data/workspace/PLL_PnL_2025A_vs_2026B_Exec_Review.pptx")


# ---- Numbers (kept explicit for exec clarity; derived from the report) ----
NUM = {
    "rev_2026_bn": 94.12,
    "rev_2025_bn": 97.41,
    "rev_yoy": -3.37,

    "ebitdar_2026_bn": 43.33,
    "ebitdar_2025_bn": 44.20,
    "ebitdar_yoy": -1.96,
    "ebitdar_m_2026": 46.0,
    "ebitdar_m_2025": 45.4,

    "ebitdar_da_2026_bn": 40.71,
    "ebitdar_da_2025_bn": 41.77,
    "ebitdar_da_yoy": -2.5,
    "ebitdar_da_m_2026": 43.25,
    "ebitdar_da_m_2025": 42.88,

    "rent_2026_bn": 36.85,
    "rent_2025_bn": 30.74,
    "rent_yoy": 19.89,
    "rent_pct_2026": 39.15,
    "rent_pct_2025": 31.56,

    "ebitda_2026_bn": 6.48,
    "ebitda_2025_bn": 13.46,
    "ebitda_yoy": -51.86,
    "ebitda_m_2026": 6.88,
    "ebitda_m_2025": 13.82,

    "npat_2026_bn": 3.86,
    "npat_2025_bn": 11.03,
    "npat_yoy": -65.03,

    "mix_2026": {
        "Nhà Liên Hoàn": 48.90,
        "Ohayo": 32.21,
        "Game-xèng": 13.01,
    },
    "mix_2025": {
        "Nhà Liên Hoàn": 48.98,
        "Ohayo": 33.65,
        "Game-xèng": 14.78,
    },

    "capex_total_bn": 2.964,
    "capex_items": [
        ("QYE", 0.653),
        ("KMO", 0.581),
        ("DLU", 0.357),
        ("GTH", 0.330),
        ("PXU", 0.264),
        ("PYE", 0.241),
        ("TNG", 0.209),
        ("HNA", 0.172),
        ("PYN", 0.158),
    ],

    "peaks": ["Apr ~10.16", "May ~12.38", "Jun ~10.54"],
    "weak": ["Oct ~6.01", "Nov ~4.85", "Dec ~4.57"],
}


# ---- Styling helpers ----

DARK = RGBColor(20, 20, 20)
MUTED = RGBColor(90, 90, 90)
ACCENT = RGBColor(0, 84, 166)  # blue
GOOD = RGBColor(0, 128, 0)
BAD = RGBColor(190, 0, 0)


def set_run_font(run, size=18, bold=False, color=DARK):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    # Title
    t = slide.shapes.title
    t.text = title
    p = t.text_frame.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = DARK

    if subtitle is not None:
        box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(12.0), Inches(0.8))
        tf = box.text_frame
        tf.clear()
        pp = tf.paragraphs[0]
        pp.text = subtitle
        pp.font.size = Pt(18)
        pp.font.color.rgb = MUTED


def add_bullets(slide, x, y, w, h, bullets, font_size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = b
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK
    return box


def add_kpi_row(slide, x, y, labels_vals, col_w=4.2):
    """labels_vals: list of (label, value_str, color_rgb)"""
    for i, (lab, val, col) in enumerate(labels_vals):
        left = Inches(x + i * col_w)
        top = Inches(y)
        width = Inches(col_w - 0.2)
        height = Inches(1.3)
        box = slide.shapes.add_textbox(left, top, width, height)
        tf = box.text_frame
        tf.clear()
        p1 = tf.paragraphs[0]
        p1.text = lab
        p1.font.size = Pt(14)
        p1.font.color.rgb = MUTED
        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(28)
        p2.font.bold = True
        p2.font.color.rgb = col


def fmt_yoy(pct: float) -> str:
    sign = "+" if pct > 0 else ""
    return f"{sign}{pct:.1f}%"


def main():
    prs = Presentation()

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    today = datetime.now(timezone.utc).strftime("%Y-%m-%d")
    add_title(slide, "PlayLand (PLL) — 2025A Results & 2026B Plan", f"Executive review deck | Source: Budget 2026_HÀ LINH.xlsx | Generated {today}")

    # --- Slide 2: Executive summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Executive summary", "Key takeaways (focus on EBITDAR & EBITDAR − D&A as operational truth)")

    add_kpi_row(
        slide,
        x=0.8,
        y=2.1,
        labels_vals=[
            ("2026 Sales budget (VND bn)", f"{NUM['rev_2026_bn']:.1f}  ({fmt_yoy(NUM['rev_yoy'])} YoY)", ACCENT),
            ("EBITDAR (VND bn)", f"{NUM['ebitdar_2026_bn']:.2f}  ({fmt_yoy(NUM['ebitdar_yoy'])} YoY)", ACCENT),
            ("EBITDAR − D&A (VND bn)", f"{NUM['ebitdar_da_2026_bn']:.2f}  ({fmt_yoy(NUM['ebitdar_da_yoy'])} YoY)", ACCENT),
        ],
        col_w=4.2,
    )

    bullets = [
        f"Core operations stable: EBITDAR margin {NUM['ebitdar_m_2026']:.1f}% vs {NUM['ebitdar_m_2025']:.1f}% (2025A).",
        f"Accounting EBITDA/NPAT drop driven by RENT treatment/level: Rent {NUM['rent_2026_bn']:.2f} bn (+{NUM['rent_yoy']:.1f}% YoY), {NUM['rent_pct_2026']:.1f}% of sales.",
        "Recommendation: use EBITDAR and EBITDAR − D&A as primary KPI set; treat EBITDA/NPAT as accounting view affected by tax optimization.",
    ]
    add_bullets(slide, 0.9, 3.7, 12.2, 2.5, bullets, font_size=18)

    # --- Slide 3: 2025 highlights & lowlights ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2025A — highlights & lowlights")

    left = [
        f"Sales: {NUM['rev_2025_bn']:.1f} bn VND",
        f"EBITDAR: {NUM['ebitdar_2025_bn']:.2f} bn (margin {NUM['ebitdar_m_2025']:.1f}%)",
        f"EBITDA: {NUM['ebitda_2025_bn']:.2f} bn (margin {NUM['ebitda_m_2025']:.1f}%)",
        f"NPAT: {NUM['npat_2025_bn']:.2f} bn",
    ]
    right = [
        f"Rent already heavy: {NUM['rent_pct_2025']:.1f}% of sales",
        "Category sensitivity: Game-xèng + Ohayo are meaningful contributors",
        "Key risk: profitability is highly sensitive to fixed cost structure (rent)",
    ]

    add_bullets(slide, 0.9, 2.0, 6.1, 4.8, ["Highlights"] + left, font_size=18)
    add_bullets(slide, 7.0, 2.0, 6.1, 4.8, ["Lowlights / risks"] + right, font_size=18)

    # --- Slide 4: 2026 plan vs 2025 (P&L lenses) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2026B vs 2025A — plan overview")

    bullets = [
        f"Sales budget: {NUM['rev_2026_bn']:.2f} bn ({fmt_yoy(NUM['rev_yoy'])} YoY).",
        f"EBITDAR: {NUM['ebitdar_2026_bn']:.2f} bn ({fmt_yoy(NUM['ebitdar_yoy'])} YoY), margin {NUM['ebitdar_m_2026']:.1f}%.",
        f"EBITDAR − D&A: {NUM['ebitdar_da_2026_bn']:.2f} bn ({fmt_yoy(NUM['ebitdar_da_yoy'])} YoY), margin {NUM['ebitdar_da_m_2026']:.2f}%.",
        f"Accounting EBITDA: {NUM['ebitda_2026_bn']:.2f} bn ({fmt_yoy(NUM['ebitda_yoy'])} YoY) driven by Rent increase to {NUM['rent_pct_2026']:.1f}% of sales.",
    ]
    add_bullets(slide, 0.9, 2.1, 12.2, 4.4, bullets, font_size=20)

    # --- Slide 5: Revenue mix ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Revenue mix (full-year)")

    mix26 = NUM["mix_2026"]
    mix25 = NUM["mix_2025"]
    total26 = sum(mix26.values())
    total25 = sum(mix25.values())

    rows = [("Category", "2025A (bn)", "2026B (bn)", "2026 share")]
    for k in ["Nhà Liên Hoàn", "Ohayo", "Game-xèng"]:
        rows.append((
            k,
            f"{mix25[k]:.2f}",
            f"{mix26[k]:.2f}",
            f"{(mix26[k]/total26*100):.1f}%",
        ))

    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(0.9), Inches(2.0), Inches(12.2), Inches(2.0)).table
    for c, h in enumerate(rows[0]):
        table.cell(0, c).text = h
    for r in range(1, len(rows)):
        for c in range(len(rows[0])):
            table.cell(r, c).text = str(rows[r][c])

    # formatting
    for r in range(len(rows)):
        for c in range(len(rows[0])):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = DARK
                    if r == 0:
                        run.font.bold = True

    add_bullets(
        slide,
        0.9,
        4.2,
        12.2,
        2.2,
        [
            "2026 plan slightly shifts mix toward Nhà Liên Hoàn (more resilient).",
            "Biggest softness vs 2025: Game-xèng and Ohayo.",
        ],
        font_size=18,
    )

    # --- Slide 6: Rent effect (why EBITDA drops) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Why EBITDA/NPAT drop (accounting view)")

    add_kpi_row(
        slide,
        x=0.8,
        y=2.1,
        labels_vals=[
            ("Rent (bn)", f"{NUM['rent_2026_bn']:.2f}  (+{NUM['rent_yoy']:.1f}% YoY)", BAD),
            ("EBITDA (bn)", f"{NUM['ebitda_2026_bn']:.2f}  ({fmt_yoy(NUM['ebitda_yoy'])} YoY)", BAD),
            ("NPAT (bn)", f"{NUM['npat_2026_bn']:.2f}  ({fmt_yoy(NUM['npat_yoy'])} YoY)", BAD),
        ],
        col_w=4.2,
    )

    add_bullets(
        slide,
        0.9,
        3.7,
        12.2,
        2.8,
        [
            f"Rent as % sales: {NUM['rent_pct_2026']:.1f}% vs {NUM['rent_pct_2025']:.1f}% (2025A).",
            "Per management note: rent is increased for tax optimization; therefore EBITDA/NPAT are not the best operational KPI.",
            "Use EBITDAR and EBITDAR − D&A to track operating performance consistently.",
        ],
        font_size=19,
    )

    # --- Slide 7: Seasonality ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2026 sales seasonality")

    add_bullets(
        slide,
        0.9,
        2.0,
        12.2,
        4.8,
        [
            "Peak months (plan): " + ", ".join(NUM["peaks"]) + " (VND bn)",
            "Weak months (plan): " + ", ".join(NUM["weak"]) + " (VND bn)",
            "Implication: staffing & marketing should align with late-spring / early-summer demand peak.",
        ],
        font_size=22,
    )

    # --- Slide 8: CAPEX summary ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "2026 CAPEX budget summary")

    add_bullets(
        slide,
        0.9,
        2.0,
        12.2,
        1.0,
        [f"Total CAPEX budget: ~{NUM['capex_total_bn']:.3f} bn VND"],
        font_size=22,
    )

    cap_rows = [("Store", "CAPEX (bn)")]
    for k, v in NUM["capex_items"]:
        cap_rows.append((k, f"{v:.3f}"))

    table = slide.shapes.add_table(len(cap_rows), 2, Inches(0.9), Inches(2.8), Inches(6.0), Inches(3.4)).table
    for c, h in enumerate(cap_rows[0]):
        table.cell(0, c).text = h
    for r in range(1, len(cap_rows)):
        table.cell(r, 0).text = cap_rows[r][0]
        table.cell(r, 1).text = cap_rows[r][1]

    for r in range(len(cap_rows)):
        for c in range(2):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(14)
                    run.font.color.rgb = DARK
                    if r == 0:
                        run.font.bold = True

    add_bullets(
        slide,
        7.1,
        2.8,
        6.0,
        3.6,
        [
            "Largest items: QYE, KMO, DLU, GTH.",
            "In workbook, monthly columns appear to be depreciation allocation (incremental D&A), not cash spend schedule.",
        ],
        font_size=18,
    )

    # --- Slide 9: KPI cheat sheet ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Key 2026 figures to remember")

    bullets = [
        f"Sales budget: {NUM['rev_2026_bn']:.2f} bn VND",
        f"EBITDAR: {NUM['ebitdar_2026_bn']:.2f} bn (margin {NUM['ebitdar_m_2026']:.1f}%)",
        f"EBITDAR − D&A: {NUM['ebitdar_da_2026_bn']:.2f} bn (margin {NUM['ebitdar_da_m_2026']:.2f}%)",
        f"Rent: {NUM['rent_2026_bn']:.2f} bn (~{NUM['rent_pct_2026']:.1f}% of sales)",
        "Primary KPI for exec cadence: EBITDAR − D&A",
    ]
    add_bullets(slide, 1.2, 2.2, 12.0, 4.5, bullets, font_size=24)

    # --- Slide 10: Appendix table (P&L) ---
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Appendix — Full-year P&L (triệu đồng)")

    pnl_rows = [
        ("Line", "2025A", "2026B", "YoY"),
        ("Revenue", "97,408.22", "94,123.48", "-3.37%"),
        ("COGS", "-23,156.54", "-22,447.99", "+3.06%"),
        ("Gross Profit", "74,251.69", "71,675.49", "-3.47%"),
        ("Total Income", "74,466.31", "72,203.28", "-3.04%"),
        ("Opex", "-30,265.01", "-28,869.36", "+4.61%"),
        ("EBITDAR", "44,201.30", "43,333.92", "-1.96%"),
        ("Rent", "-30,739.57", "-36,854.02", "+19.89%"),
        ("EBITDA", "13,461.73", "6,479.90", "-51.86%"),
        ("Depreciation", "-2,433.80", "-2,623.94", "+7.81%"),
        ("EBIT", "11,027.94", "3,855.96", "-65.03%"),
        ("NPAT", "11,027.94", "3,855.96", "-65.03%"),
    ]

    table = slide.shapes.add_table(len(pnl_rows), 4, Inches(0.7), Inches(2.0), Inches(12.6), Inches(4.6)).table
    for c in range(4):
        table.cell(0, c).text = pnl_rows[0][c]
    for r in range(1, len(pnl_rows)):
        for c in range(4):
            table.cell(r, c).text = pnl_rows[r][c]

    for r in range(len(pnl_rows)):
        for c in range(4):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = DARK
                    if r == 0:
                        run.font.bold = True
                    if pnl_rows[r][0] in ("EBITDAR", "EBITDA", "NPAT") and c == 0:
                        run.font.bold = True

    prs.save(OUT_PATH)
    print(f"Wrote: {OUT_PATH}")


if __name__ == "__main__":
    main()
