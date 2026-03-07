#!/usr/bin/env python3
"""make_pll_pnl_deck_v2.py

Create a more polished executive PPTX deck for PLL (2025A vs 2026B).

Output:
  /data/workspace/PLL_PnL_2025A_vs_2026B_Exec_Review_v2.pptx

Design goals:
- Consistent theme (dark title + light content)
- KPI cards, charts, clean spacing
- Minimal tables (only appendix)

Deps:
  pip install python-pptx
"""

from __future__ import annotations

from pathlib import Path
from datetime import datetime, timezone

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION

OUT_PATH = Path("/data/workspace/PLL_PnL_2025A_vs_2026B_Exec_Review_v2.pptx")

# Numbers (from the approved report)
NUM = {
    "rev": {"2025A": 97.408, "2026B": 94.123, "yoy": -3.37},
    "ebitdar": {"2025A": 44.201, "2026B": 43.334, "yoy": -1.96, "m25": 45.4, "m26": 46.0},
    "ebitdar_da": {"2025A": 41.768, "2026B": 40.710, "yoy": -2.5, "m25": 42.88, "m26": 43.25},
    "rent": {"2025A": 30.740, "2026B": 36.854, "yoy": 19.89, "pct25": 31.56, "pct26": 39.15},
    "ebitda": {"2025A": 13.462, "2026B": 6.480, "yoy": -51.86, "m25": 13.82, "m26": 6.88},
    "npat": {"2025A": 11.028, "2026B": 3.856, "yoy": -65.03},
    "mix_2026": {"Nhà Liên Hoàn": 48.900, "Ohayo": 32.209, "Game-xèng": 13.014},
    "capex_total": 2.964,
    "capex_items": [("QYE", 0.653), ("KMO", 0.581), ("DLU", 0.357), ("GTH", 0.330), ("PXU", 0.264)],
}

# Theme
NAVY = RGBColor(10, 26, 57)
BLUE = RGBColor(0, 84, 166)
ICE = RGBColor(235, 243, 255)
TEXT = RGBColor(25, 25, 25)
MUTED = RGBColor(90, 90, 90)
WHITE = RGBColor(255, 255, 255)
RED = RGBColor(190, 0, 0)
GREEN = RGBColor(0, 120, 0)

FONT_TITLE = "Calibri"
FONT_BODY = "Calibri"


def add_top_bar(slide, color=BLUE):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.28))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def set_text(tf, text, size=18, bold=False, color=TEXT, align=PP_ALIGN.LEFT):
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_BODY
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align


def add_title(slide, title, subtitle=None, dark=False):
    if not dark:
        add_top_bar(slide)
    left = Inches(0.8)
    top = Inches(0.55) if not dark else Inches(1.2)
    box = slide.shapes.add_textbox(left, top, Inches(11.7), Inches(0.9))
    tf = box.text_frame
    set_text(tf, title, size=38 if not dark else 44, bold=True, color=(WHITE if dark else TEXT))

    if subtitle:
        box2 = slide.shapes.add_textbox(left, top + Inches(0.8), Inches(11.7), Inches(0.6))
        tf2 = box2.text_frame
        set_text(tf2, subtitle, size=18, bold=False, color=(ICE if dark else MUTED))


def add_kpi_card(slide, x, y, w, h, label, value, note=None, accent=BLUE, bg=WHITE):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = RGBColor(220, 228, 240)

    # left accent strip
    strip = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(x), Inches(y), Inches(0.10), Inches(h))
    strip.fill.solid(); strip.fill.fore_color.rgb = accent
    strip.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.18), Inches(w - 0.3), Inches(0.4)).text_frame
    set_text(t1, label, size=13, bold=True, color=MUTED)

    t2 = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.55), Inches(w - 0.3), Inches(0.65)).text_frame
    set_text(t2, value, size=28, bold=True, color=TEXT)

    if note:
        t3 = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 1.08), Inches(w - 0.3), Inches(0.35)).text_frame
        set_text(t3, note, size=12, bold=False, color=MUTED)


def add_bullets(slide, x, y, w, h, bullets, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(size)
        p.font.color.rgb = TEXT


def add_clustered_bar(slide, x, y, w, h, title, cats, series):
    """series: list of (name, values list)."""
    chart_data = ChartData()
    chart_data.categories = cats
    for sname, vals in series:
        chart_data.add_series(sname, vals)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data,
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    chart.has_legend = True
    chart.legend.include_in_layout = False

    # data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.position = XL_LABEL_POSITION.OUTSIDE_END
    dl.number_format = "0.0"

    # colors
    for i, s in enumerate(chart.series):
        fill = s.format.fill
        fill.solid()
        fill.fore_color.rgb = [BLUE, NAVY][i % 2]

    return chart


def add_donut(slide, x, y, w, h, title, labels, values):
    chart_data = ChartData()
    chart_data.categories = labels
    chart_data.add_series("2026B", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT,
        Inches(x), Inches(y), Inches(w), Inches(h),
        chart_data,
    ).chart

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)

    chart.has_legend = True
    chart.legend.include_in_layout = False

    # data labels
    plot = chart.plots[0]
    plot.has_data_labels = True
    dl = plot.data_labels
    dl.number_format = '0%'
    dl.show_percentage = True

    palette = [RGBColor(0, 84, 166), RGBColor(10, 26, 57), RGBColor(0, 150, 136)]
    for i, pt in enumerate(chart.series[0].points):
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = palette[i % len(palette)]

    return chart


def yoy(p):
    return ("+" if p > 0 else "") + f"{p:.1f}%"


def main():
    prs = Presentation()
    # widescreen
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    today = datetime.now(timezone.utc).strftime("%d %b %Y")

    # Slide 1: dark title
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid(); bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    add_title(slide, "PlayLand (PLL)", "2025A Results & 2026B Plan — Executive Review", dark=True)
    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(12.0), Inches(0.4)).text_frame
    set_text(foot, f"Source: Budget 2026_HÀ LINH.xlsx | Generated {today}", size=12, color=ICE)

    # Slide 2: Executive summary (KPI cards)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Executive summary")

    add_kpi_card(slide, 0.9, 1.8, 3.9, 1.45, "2026 Sales budget (bn VND)", f"{NUM['rev']['2026B']:.1f}", f"YoY {yoy(NUM['rev']['yoy'])}")
    add_kpi_card(slide, 4.95, 1.8, 3.9, 1.45, "EBITDAR (bn VND)", f"{NUM['ebitdar']['2026B']:.2f}", f"YoY {yoy(NUM['ebitdar']['yoy'])}")
    add_kpi_card(slide, 9.00, 1.8, 3.9, 1.45, "EBITDAR − D&A (bn VND)", f"{NUM['ebitdar_da']['2026B']:.2f}", f"YoY {yoy(NUM['ebitdar_da']['yoy'])}")

    add_bullets(
        slide,
        0.95,
        3.55,
        12.0,
        3.2,
        [
            f"Core ops stable: EBITDAR margin {NUM['ebitdar']['m26']:.1f}% vs {NUM['ebitdar']['m25']:.1f}% (2025A).",
            f"Accounting EBITDA/NPAT drop is driven by RENT treatment/level (tax optimization): Rent {NUM['rent']['pct26']:.1f}% of sales.",
            "Primary KPI set for exec cadence: EBITDAR and EBITDAR − D&A.",
        ],
        size=20,
    )

    # Slide 3: Plan vs Actual (bars)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "2026B vs 2025A — KPI comparison")

    cats = ["Sales", "EBITDAR", "EBITDAR−D&A"]
    series = [
        ("2025A", [NUM['rev']['2025A'], NUM['ebitdar']['2025A'], NUM['ebitdar_da']['2025A']]),
        ("2026B", [NUM['rev']['2026B'], NUM['ebitdar']['2026B'], NUM['ebitdar_da']['2026B']]),
    ]
    add_clustered_bar(slide, 0.9, 2.0, 7.2, 4.6, "Absolute (bn VND)", cats, series)

    add_kpi_card(
        slide,
        8.35,
        2.0,
        4.55,
        1.25,
        "EBITDAR margin",
        f"{NUM['ebitdar']['m26']:.1f}%",
        f"2025A: {NUM['ebitdar']['m25']:.1f}%",
        accent=BLUE,
        bg=ICE,
    )
    add_kpi_card(
        slide,
        8.35,
        3.45,
        4.55,
        1.25,
        "EBITDAR − D&A margin",
        f"{NUM['ebitdar_da']['m26']:.2f}%",
        f"2025A: {NUM['ebitdar_da']['m25']:.2f}%",
        accent=BLUE,
        bg=ICE,
    )
    add_kpi_card(
        slide,
        8.35,
        4.90,
        4.55,
        1.25,
        "YoY sales",
        yoy(NUM['rev']['yoy']),
        "(headline: slightly down)",
        accent=BLUE,
        bg=ICE,
    )

    # Slide 4: Why EBITDA drops (rent)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Accounting drop explained (rent effect)")

    add_kpi_card(slide, 0.9, 1.9, 4.1, 1.4, "Rent (bn VND)", f"{NUM['rent']['2026B']:.2f}", f"YoY +{NUM['rent']['yoy']:.1f}%", accent=RED)
    add_kpi_card(slide, 5.2, 1.9, 3.8, 1.4, "EBITDA (bn VND)", f"{NUM['ebitda']['2026B']:.2f}", f"YoY {yoy(NUM['ebitda']['yoy'])}", accent=RED)
    add_kpi_card(slide, 9.2, 1.9, 3.7, 1.4, "NPAT (bn VND)", f"{NUM['npat']['2026B']:.2f}", f"YoY {yoy(NUM['npat']['yoy'])}", accent=RED)

    # small chart rent as % sales
    cats = ["2025A", "2026B"]
    series = [("Rent % Sales", [NUM['rent']['pct25'], NUM['rent']['pct26']])]
    chart_data = ChartData(); chart_data.categories = cats; chart_data.add_series("Rent%", [NUM['rent']['pct25'], NUM['rent']['pct26']])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.9), Inches(3.7), Inches(6.2), Inches(3.0), chart_data).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "Rent as % of sales"
    chart.plots[0].has_data_labels = True
    dl = chart.plots[0].data_labels
    dl.number_format = "0.0"; dl.position = XL_LABEL_POSITION.OUTSIDE_END
    chart.series[0].format.fill.solid(); chart.series[0].format.fill.fore_color.rgb = RED

    add_bullets(
        slide,
        7.35,
        3.75,
        5.6,
        3.0,
        [
            "Management note: rent is increased to optimize tax.",
            "Therefore EBITDA/NPAT are not the best operational KPI.",
            "Use EBITDAR / EBITDAR − D&A for performance tracking.",
        ],
        size=18,
    )

    # Slide 5: Revenue mix donut
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "2026B revenue mix")

    labels = list(NUM["mix_2026"].keys())
    values = list(NUM["mix_2026"].values())
    add_donut(slide, 0.9, 2.0, 6.2, 4.8, "Share of 2026B sales", labels, values)

    add_bullets(
        slide,
        7.35,
        2.1,
        5.6,
        4.8,
        [
            "Nhà Liên Hoàn is the largest share and planned to be stable.",
            "Softness vs 2025 mainly comes from Game-xèng and Ohayo.",
            "Implication: protect Ohayo & Game-xèng demand while leveraging NLH base.",
        ],
        size=18,
    )

    # Slide 6: CAPEX
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "2026 CAPEX plan")

    add_kpi_card(slide, 0.9, 1.9, 4.6, 1.35, "Total CAPEX (bn VND)", f"{NUM['capex_total']:.3f}", "Top items shown right", accent=GREEN)

    # top-5 bar
    stores = [k for k, _ in NUM["capex_items"]]
    vals = [v for _, v in NUM["capex_items"]]
    cd = ChartData(); cd.categories = stores; cd.add_series("CAPEX", vals)
    ch = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Inches(5.8), Inches(1.9), Inches(7.0), Inches(3.7), cd).chart
    ch.has_title = True
    ch.chart_title.text_frame.text = "Top CAPEX stores (bn VND)"
    ch.plots[0].has_data_labels = True
    ch.plots[0].data_labels.number_format = "0.000"

    add_bullets(
        slide,
        0.9,
        3.45,
        4.6,
        3.9,
        [
            "Largest CAPEX: QYE, KMO, DLU, GTH, PXU.",
            "Workbook indicates monthly columns may represent depreciation allocation (incremental D&A) rather than cash schedule.",
        ],
        size=16,
    )

    # Slide 7: Appendix (compact table)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Appendix — full-year P&L (triệu đồng)")

    pnl = [
        ("Revenue", 97408.22, 94123.48, -3.37),
        ("COGS", -23156.54, -22447.99, 3.06),
        ("Gross Profit", 74251.69, 71675.49, -3.47),
        ("Opex", -30265.01, -28869.36, 4.61),
        ("EBITDAR", 44201.30, 43333.92, -1.96),
        ("Rent", -30739.57, -36854.02, 19.89),
        ("EBITDA", 13461.73, 6479.90, -51.86),
        ("Depreciation", -2433.80, -2623.94, 7.81),
        ("NPAT", 11027.94, 3855.96, -65.03),
    ]

    rows = len(pnl) + 1
    cols = 4
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.9), Inches(12.4), Inches(5.2)).table
    headers = ["Line", "2025A", "2026B", "YoY"]
    for c, h in enumerate(headers):
        table.cell(0, c).text = h

    for r, (line, a25, a26, py) in enumerate(pnl, start=1):
        table.cell(r, 0).text = line
        table.cell(r, 1).text = f"{a25:,.2f}"
        table.cell(r, 2).text = f"{a26:,.2f}"
        table.cell(r, 3).text = yoy(py)

    # style header
    for c in range(cols):
        cell = table.cell(0, c)
        cell.fill.solid(); cell.fill.fore_color.rgb = ICE
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(12)

    # style body
    for r in range(1, rows):
        for c in range(cols):
            tf = table.cell(r, c).text_frame
            for p in tf.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(12)
                    run.font.color.rgb = TEXT

    prs.save(OUT_PATH)
    print(f"Wrote: {OUT_PATH}")


if __name__ == "__main__":
    main()
